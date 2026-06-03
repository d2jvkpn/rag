#!/usr/bin/env python3

import json
import os
import re
import sys
from pathlib import Path


def fail(message: str) -> int:
    print(message, file=sys.stderr)
    return 1


def clean_cell(cell) -> str:
    if cell is None:
        return ""
    text = str(cell).replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    text = "<br>".join(lines)
    text = re.sub(r"\s+", " ", text).strip()
    return text.replace("|", r"\|")


def format_markdown_table(rows) -> str:
    normalized = []
    max_cols = 0
    for row in rows:
        if row is None:
            continue
        cells = [clean_cell(cell) for cell in row]
        if not any(cells):
            continue
        normalized.append(cells)
        max_cols = max(max_cols, len(cells))

    if not normalized or max_cols == 0:
        return ""

    for row in normalized:
        if len(row) < max_cols:
            row.extend([""] * (max_cols - len(row)))

    header = normalized[0]
    if not any(header):
        header = [f"列{i + 1}" for i in range(max_cols)]
        normalized[0] = header

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * max_cols) + " |",
    ]
    for row in normalized[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _pptx_image_label_url(shape):
    """Return (label, external_url) from <p:cNvPr descr/name> attributes.

    If descr is an external URL it becomes the ref URL instead of the label,
    matching the Go parser's behaviour.
    """
    try:
        from pptx.oxml.ns import qn
        cNvPr = shape.element.find(".//" + qn("p:cNvPr"))
        if cNvPr is not None:
            descr = cNvPr.get("descr") or ""
            if descr:
                if descr.startswith("http://") or descr.startswith("https://"):
                    return "", descr
                return descr, ""
            name = cNvPr.get("name") or ""
            return name, ""
    except Exception:
        pass
    return "", ""


def extract_slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if not shape.has_text_frame or shape.has_table:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def extract_slide_tables(slide):
    tables = []
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        rows = [[cell.text for cell in row.cells] for row in tbl.rows]
        if format_markdown_table(rows):
            tables.append(rows)
    return tables


def extract_slide_links(slide, slide_number: int):
    links = []
    index = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    url = run.hyperlink.address
                except Exception:
                    url = None
                if not url:
                    continue
                index += 1
                anchor = run.text.strip()
                links.append({
                    "ref_id": f"pptx-link-{slide_number}-{index}",
                    "ref_type": "link",
                    "label": anchor or url,
                    "anchor_text": anchor,
                    "page": slide_number,
                    "url": url,
                    "is_external": True,
                })
    return links


def save_slide_image(blob: bytes, ref_id: str, media_dir, storage_base) -> str:
    if media_dir is None or storage_base is None:
        return ""
    try:
        media_dir.mkdir(parents=True, exist_ok=True)
        filename = f"{ref_id}.png"
        disk_path = media_dir / filename
        disk_path.write_bytes(blob)
        return str(disk_path.relative_to(storage_base))
    except Exception:
        return ""


def extract_slide_images(slide, slide_number: int, media_dir=None, storage_base=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    refs = []
    image_index = 0
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            blob = shape.image.blob
        except Exception:
            continue
        image_index += 1
        label, external_url = _pptx_image_label_url(shape)
        ref_id = f"pptx-image-{slide_number}-{image_index}"
        ref = {
            "ref_id": ref_id,
            "ref_type": "image",
            "label": label,
            "caption": label,
            "page": slide_number,
        }
        if external_url:
            ref["url"] = external_url
            ref["is_external"] = True
        storage_path = save_slide_image(blob, ref_id, media_dir, storage_base)
        if storage_path:
            ref["storage_path"] = storage_path
        refs.append(ref)
    return refs


def extract_slide_notes(slide) -> str:
    try:
        if slide.has_notes_slide:
            return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        pass
    return ""


def render_slide(slide_number: int, text: str, tables, image_refs, note: str) -> str:
    """Render slide to text string, matching Go's output structure:

    - ``Slide N`` label as first line
    - body text and Markdown tables
    - image placeholders (``[Image:ref_id label?]``) appended after body
    - speaker notes appended last as ``Notes: …``
    """
    parts = [f"Slide {slide_number}"]
    if text:
        parts.append(text)
    for rows in tables:
        md = format_markdown_table(rows)
        if md:
            parts.append(md)
    for ref in image_refs:
        label = ref.get("label") or ""
        ref_id = ref["ref_id"]
        placeholder = f"[Image:{ref_id} {label}]" if label else f"[Image:{ref_id}]"
        parts.append(placeholder)
    if note:
        parts.append(f"Notes: {note}")
    return "\n".join(p for p in parts if p.strip()).strip()


def main() -> int:
    if len(sys.argv) not in (2, 3):
        return fail("usage: parse_pptx.py <pptx_path> [media_dir]")

    path = Path(sys.argv[1])
    if not path.is_file():
        return fail(f"pptx file not found: {path}")

    media_dir = Path(sys.argv[2]) if len(sys.argv) == 3 and sys.argv[2] else None
    storage_base_env = os.environ.get("PPTX_PARSER_STORAGE_BASE")
    storage_base = Path(storage_base_env) if storage_base_env else (
        media_dir.parent if media_dir is not None else None
    )

    try:
        from pptx import Presentation
    except Exception as exc:
        return fail(f"python-pptx import failed: {exc}")

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return fail(f"open pptx failed: {exc}")

    total_slides = len(prs.slides)
    page_data = []
    failed_slides = 0

    for index, slide in enumerate(prs.slides, start=1):
        try:
            text = extract_slide_text(slide)
            tables = extract_slide_tables(slide)
            image_refs = extract_slide_images(slide, index, media_dir, storage_base)
            link_refs = extract_slide_links(slide, index)
            note = extract_slide_notes(slide)
        except Exception as exc:
            failed_slides += 1
            if failed_slides == total_slides:
                return fail(f"pptx extraction failed: {exc}")
            continue

        has_content = bool(text or tables or image_refs or link_refs or note)
        if not has_content:
            continue

        rendered = render_slide(index, text, tables, image_refs, note)
        all_refs = image_refs + link_refs

        table_refs = []
        for ti, rows in enumerate(tables, start=1):
            label = f"Table {ti} (Slide {index})"
            table_refs.append({
                "ref_id": f"pptx-table-{index}-{ti}",
                "ref_type": "table",
                "label": label,
                "caption": label,
                "page": index,
            })

        refs = all_refs + table_refs
        page_data.append({"text": rendered, "page": index, "refs": refs})

    if not page_data:
        return fail("pptx extraction failed: no extractable content found")

    payload = {
        "text": "\n\n".join(p["text"] for p in page_data).strip(),
        "page_count": len(page_data),
        "pages": page_data,
    }
    json.dump(payload, sys.stdout, ensure_ascii=False)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
