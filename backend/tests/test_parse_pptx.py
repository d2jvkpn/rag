import importlib.util
import pathlib
import unittest


SCRIPT_PATH = pathlib.Path(__file__).parent.parent / "scripts" / "parse_pptx.py"
SPEC = importlib.util.spec_from_file_location("parse_pptx", SCRIPT_PATH)
parse_pptx = importlib.util.module_from_spec(SPEC)
assert SPEC.loader is not None
SPEC.loader.exec_module(parse_pptx)


# ---------------------------------------------------------------------------
# Minimal fake objects — just enough to exercise each extraction function.
# ---------------------------------------------------------------------------

class FakeHyperlink:
    def __init__(self, address):
        self.address = address


class FakeRun:
    def __init__(self, text, url=None):
        self.text = text
        self.hyperlink = FakeHyperlink(url) if url else FakeHyperlink(None)


class FakeParagraph:
    def __init__(self, text, runs=None):
        self._text = text
        self.runs = runs or [FakeRun(text)]

    @property
    def text(self):
        return self._text


class FakeTextFrame:
    def __init__(self, paragraphs):
        self.paragraphs = paragraphs


class FakeCell:
    def __init__(self, text):
        self.text = text


class FakeRow:
    def __init__(self, cells):
        self.cells = [FakeCell(c) for c in cells]


class FakeTable:
    def __init__(self, rows):
        self.rows = [FakeRow(r) for r in rows]


class FakeImage:
    def __init__(self, blob=b"png"):
        self.blob = blob


class FakeElement:
    def find(self, *args, **kwargs):
        return None


PICTURE = 13
TEXTBOX = 17


class FakeShape:
    def __init__(self, shape_type=TEXTBOX, text_frame=None, table=None, image=None):
        self.shape_type = shape_type
        self._text_frame = text_frame
        self._table = table
        self._image = image
        self.element = FakeElement()
        # shape.name for image shapes
        self.name = ""

    @property
    def has_text_frame(self):
        return self._text_frame is not None

    @property
    def has_table(self):
        return self._table is not None

    @property
    def text_frame(self):
        return self._text_frame

    @property
    def table(self):
        return self._table

    @property
    def image(self):
        return self._image


class FakeNotesTextFrame:
    def __init__(self, text):
        self.text = text


class FakeNotesSlide:
    def __init__(self, text):
        self.notes_text_frame = FakeNotesTextFrame(text)


class FakeSlide:
    def __init__(self, shapes, notes_text=""):
        self.shapes = shapes
        self._notes_text = notes_text

    @property
    def has_notes_slide(self):
        return bool(self._notes_text)

    @property
    def notes_slide(self):
        return FakeNotesSlide(self._notes_text)


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

class ParsePPTXTextTests(unittest.TestCase):
    def test_extract_text_from_text_frames(self):
        slide = FakeSlide([
            FakeShape(shape_type=TEXTBOX, text_frame=FakeTextFrame([
                FakeParagraph("Title"),
                FakeParagraph("Body line"),
            ])),
        ])
        text = parse_pptx.extract_slide_text(slide)
        self.assertIn("Title", text)
        self.assertIn("Body line", text)

    def test_text_frame_with_table_skipped(self):
        tf = FakeTextFrame([FakeParagraph("should be skipped")])
        slide = FakeSlide([
            FakeShape(shape_type=TEXTBOX, text_frame=tf, table=FakeTable([["A", "B"]])),
        ])
        text = parse_pptx.extract_slide_text(slide)
        self.assertEqual(text, "")


class ParsePPTXTableTests(unittest.TestCase):
    def test_table_extracted(self):
        slide = FakeSlide([
            FakeShape(shape_type=TEXTBOX, table=FakeTable([["Name", "Score"], ["Alice", "90"]])),
        ])
        tables = parse_pptx.extract_slide_tables(slide)
        self.assertEqual(len(tables), 1)
        md = parse_pptx.format_markdown_table(tables[0])
        self.assertIn("| Name | Score |", md)

    def test_table_rendered_without_label_prefix(self):
        slide = FakeSlide([
            FakeShape(shape_type=TEXTBOX, text_frame=FakeTextFrame([FakeParagraph("Intro")])),
            FakeShape(shape_type=TEXTBOX, table=FakeTable([["A", "B"], ["1", "2"]])),
        ])
        text = parse_pptx.extract_slide_text(slide)
        tables = parse_pptx.extract_slide_tables(slide)
        rendered = parse_pptx.render_slide(1, text, tables, [], "")
        self.assertIn("| A | B |", rendered)
        self.assertNotIn("Table 1 (Slide", rendered)


class ParsePPTXImageTests(unittest.TestCase):
    def test_no_image_for_non_picture_shape(self):
        slide = FakeSlide([
            FakeShape(shape_type=TEXTBOX, text_frame=FakeTextFrame([FakeParagraph("text")])),
        ])
        refs = parse_pptx.extract_slide_images(slide, 1)
        self.assertEqual(refs, [])

    def test_image_ref_for_picture_shape(self):
        slide = FakeSlide([
            FakeShape(shape_type=PICTURE, image=FakeImage(b"fakeblob")),
        ])
        refs = parse_pptx.extract_slide_images(slide, 3)
        self.assertEqual(len(refs), 1)
        self.assertEqual(refs[0]["ref_type"], "image")
        self.assertEqual(refs[0]["ref_id"], "pptx-image-3-1")
        self.assertNotIn("storage_path", refs[0])

    def test_image_placeholder_in_rendered_text(self):
        slide = FakeSlide([
            FakeShape(shape_type=PICTURE, image=FakeImage(b"x")),
        ])
        image_refs = parse_pptx.extract_slide_images(slide, 2)
        rendered = parse_pptx.render_slide(2, "", [], image_refs, "")
        self.assertIn("[Image:pptx-image-2-1]", rendered)

    def test_image_placeholder_includes_label(self):
        class ShapeWithLabel(FakeShape):
            pass

        shape = ShapeWithLabel(shape_type=PICTURE, image=FakeImage(b"x"))
        # Override _pptx_image_label_url to return a known label
        import unittest.mock as mock
        with mock.patch.object(parse_pptx, "_pptx_image_label_url", return_value=("MyFigure", "")):
            slide = FakeSlide([shape])
            refs = parse_pptx.extract_slide_images(slide, 1)
            rendered = parse_pptx.render_slide(1, "", [], refs, "")
            self.assertIn("[Image:pptx-image-1-1 MyFigure]", rendered)

    def test_image_with_external_url_stored_in_ref(self):
        import unittest.mock as mock
        shape = FakeShape(shape_type=PICTURE, image=FakeImage(b"x"))
        with mock.patch.object(parse_pptx, "_pptx_image_label_url",
                               return_value=("", "https://ext.example.com/img.png")):
            slide = FakeSlide([shape])
            refs = parse_pptx.extract_slide_images(slide, 1)
            self.assertEqual(refs[0].get("url"), "https://ext.example.com/img.png")
            self.assertTrue(refs[0].get("is_external"))

    def test_image_saved_to_storage(self):
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            storage_base = pathlib.Path(tmp)
            media_dir = storage_base / "doc-1"
            slide = FakeSlide([FakeShape(shape_type=PICTURE, image=FakeImage(b"PNG"))])
            refs = parse_pptx.extract_slide_images(slide, 2, media_dir, storage_base)
            self.assertIn("storage_path", refs[0])
            self.assertTrue((storage_base / refs[0]["storage_path"]).is_file())


class ParsePPTXNotesTests(unittest.TestCase):
    def test_notes_appended_to_rendered_text(self):
        slide = FakeSlide(
            [FakeShape(shape_type=TEXTBOX, text_frame=FakeTextFrame([FakeParagraph("Body")]))],
            notes_text="Speaker reminder",
        )
        note = parse_pptx.extract_slide_notes(slide)
        self.assertEqual(note, "Speaker reminder")
        rendered = parse_pptx.render_slide(1, "Body", [], [], note)
        self.assertIn("Notes: Speaker reminder", rendered)

    def test_no_notes_when_slide_has_none(self):
        slide = FakeSlide([])
        note = parse_pptx.extract_slide_notes(slide)
        self.assertEqual(note, "")


class ParsePPTXRenderTests(unittest.TestCase):
    def test_slide_n_prefix_in_rendered_text(self):
        rendered = parse_pptx.render_slide(5, "Content", [], [], "")
        self.assertTrue(rendered.startswith("Slide 5"))

    def test_render_order_body_images_notes(self):
        image_refs = [{"ref_id": "pptx-image-1-1", "label": ""}]
        rendered = parse_pptx.render_slide(1, "Body", [], image_refs, "Notes here")
        body_pos = rendered.index("Body")
        image_pos = rendered.index("[Image:")
        notes_pos = rendered.index("Notes:")
        self.assertLess(body_pos, image_pos)
        self.assertLess(image_pos, notes_pos)

    def test_page_count_equals_non_empty_slides(self):
        """page_count in payload = number of slides with content, not total slides."""
        import json, io, sys, tempfile
        from pptx import Presentation
        from pptx.util import Inches
        from pathlib import Path

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # completely blank
        content_layout = prs.slide_layouts[5]  # blank (can add shapes)
        # Slide 1: has text
        s1 = prs.slides.add_slide(content_layout)
        tf = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf.text_frame.text = "Hello"
        # Slide 2: empty
        prs.slides.add_slide(blank_layout)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            tmp = f.name

        old_argv, old_stdout = sys.argv, sys.stdout
        try:
            sys.argv = ["parse_pptx.py", tmp]
            buf = io.StringIO()
            sys.stdout = buf
            parse_pptx.main()
            sys.stdout = old_stdout
            payload = json.loads(buf.getvalue())
            # Only non-empty slides counted
            self.assertLess(payload["page_count"], 2)
        finally:
            sys.argv = old_argv
            sys.stdout = old_stdout
            Path(tmp).unlink(missing_ok=True)

    def test_link_extraction(self):
        slide = FakeSlide([
            FakeShape(shape_type=TEXTBOX, text_frame=FakeTextFrame([
                FakeParagraph("click", runs=[FakeRun("click", url="https://example.com")]),
            ])),
        ])
        links = parse_pptx.extract_slide_links(slide, 1)
        self.assertEqual(len(links), 1)
        self.assertEqual(links[0]["url"], "https://example.com")
        self.assertEqual(links[0]["ref_id"], "pptx-link-1-1")


if __name__ == "__main__":
    unittest.main()
